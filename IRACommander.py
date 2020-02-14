#! /usr/bin/env python3
# -*- coding: utf-8 -*-
from tkinter import *
from tkinter import PhotoImage, Label, Frame
from tkinter.filedialog import askopenfilename
from audioanalysis import audio2text
import os
from pptx import Presentation
import csv
from arduserial import arduino_port
import pyttsx3
from fuzzywuzzy import fuzz, process
from time import sleep


class IRACommander(Frame):
    def __init__(self, master=None):
        Frame.__init__(self, master)
        self.master.title("IRA Commander")
        self.master.iconbitmap('icon.ico')
        self.master.minsize(width=260, height=140)
        self.master.maxsize(width=260, height=140)
        self.grid()

    def quit(self):
        ira.stop()
        arduino.close()
        self.destroy()
        sys.exit()


class CoreGUI(object):
    def __init__(self, parent):
        self.fileButton = Button(parent, text="Presentation")
        self.fileButton.bind("<Button-1>", file_input)
        self.fileButton.pack(fill=X, padx=90, pady=10)
        self.audioButton = Button(parent, text="Conversation")
        self.audioButton.bind("<Button-1>", audio_input)
        self.audioButton.pack(fill=X, padx=90, pady=10)
        self.brandButton = Button(parent, text="Branding")
        self.brandButton.bind("<Button-1>", branding)
        self.brandButton.pack(fill=X, padx=90, pady=10)

    def button_state(self):
        self.fileButton.config(relief=RAISED)


def file_input(event):
    command = ""
    name = askopenfilename(initialdir="/",
                           filetypes=(("Text File", "*.txt"),
                                      ("CSV File", "*.csv"),
                                      ("Presentation File", "*.pptx"),
                                      ("All Files", "*.*")),
                           title="Choose a file."
                           )
    filename, ext = os.path.splitext(name)
    try:
        file = open(name, 'r', encoding="utf8")
        if ext == ".pptx":
            text_runs = []
            prs = Presentation(name)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        stripped = paragraph.text.strip()
                        if stripped:
                            text_runs.append(paragraph.text)
            command = text_runs[0]
        elif ext == ".csv":
            csv_data = csv.reader(file)
            for data in csv_data:
                word = 0
                command += data[word]
                command += " "
                word += 1
            command = command.rstrip()
            file.close()
        else:
            command = file.read()
            file.close()
        presentation_exec(command)
        tk.button_state()
    except FileNotFoundError:
        print("No file choosen")
        tk.button_state()


def presentation_exec(command):
    say("I will start whenever you say, start")
    ira.runAndWait()
    while True:
        start = audio2text()
        start_command = str(start).lower()
        print(start_command)
        split_command = start_command.split()
        if "start" or "star" in split_command:
            say(command)
            # say("Thank you for Listening. Have a nice day.")
            ira.runAndWait()
            break
        else:
            pass
    return


def audio_input(event):
    silent = 0
    ara = ["ara", "ira", "eira", "aera", "ayra","hira", "aa raha", "hara", "haha", "aaram", "aara", "aarav", "vera", "yaara",
           "yara", "are"]
    off = ["shut down", "turn off", "turner", "shutdown", "turnoff"]
    hello = ["hello", "hello ira", "hi ira", "hi", "hello ara", "hi ara", "hello hira", "hi hira",
             "hirer", "hi there"]
    ambiguous = ["what", "is", "the", "what is the", "when is the", "when", "is the", "do", "you", "do you",
                 "what is", "when is"]
    bye = ["by", "bye", "goodbye", "good bye", "allah hafez", "allah hafiz", "ta ta", "tata"]
    silent_mode = True
    while silent_mode:
        command = audio2text()
        call = str(command).lower()
        words_in_call = call.split()
        find_ira = set(words_in_call).intersection(set(ara))
        print(call)
        if find_ira:
            say("Yes,")
            ira.runAndWait()
            silent_mode = False
            while True:
                command = audio2text()
                answer = str(command).lower()
                # print("answer passed")
                if not answer:
                    # print("not answer passed")
                    silent += 1
                    if silent < 5:
                        print(silent)
                        pass
                    else:
                        silent = 0
                        say("I am going to silent mode. Call me, whenever you want to back in conversation")
                        ira.runAndWait()
                        silent_mode = True
                        break
                else:
                    if silent_mode is False:
                        silent = 0
                        bye_id = [bye.index(i) for j in answer.split() for i in bye if j in i]
                        match, match_percent, dictionary, ans4match = match_counter(answer)
                        # print("answer condition passed")
                        if answer in ara:
                            say("yes?")
                        elif answer in hello:
                            # print("answer speech condition passed")
                            say("Hi there, how are you?")
                            # print("answer speech passed")
                        elif answer.find("i am fine") != -1 or answer.find("i'm fine") != -1:
                            say("I am glad to know")
                        elif answer.find("what is your name") != -1:
                            say("My name is eira")
                        elif answer.find("how are you") != -1:
                            say("thank you for asking, I am fine.")
                        elif answer.find("thank you") != -1:
                            say("You're Welcome")
                        elif answer.find("say hi to") != -1 or answer.find("say hello to") != -1:
                            hi_command = answer.split()
                            greet_type = hi_command.index("say") + 1
                            consignee = hi_command.index("to") + 1
                            try:
                                say(hi_command[greet_type] + " " + hi_command[consignee] + " " + hi_command[
                                    consignee + 1] + ", how are you?")
                            except IndexError:
                                try:
                                    say(hi_command[greet_type] + " " + hi_command[consignee] + ", how are you?")
                                except IndexError:
                                    say("I think, you didn't completed your sentence")
                        elif match_percent >= 70 and answer not in ambiguous:
                            say(dictionary[ans4match])
                        elif bye_id:
                            say(bye[bye_id[0]])
                            say("It was a nice conversation with you, Thank you")
                            ira.runAndWait()
                            silent_mode = True
                            break
                        elif answer.find("close the program") != -1 or answer in off:
                            say("Terminating the eira Commander, have a nice day")
                            ira.runAndWait()
                            root.quit()
                        else:
                            say("I heard")
                            say(answer)
                            say("I didn't learnt any replies for that.")
                        ira.runAndWait()
                        print(answer)
        else:
            continue
    return


def branding(event):
    verses = [line.rstrip('\n') for line in open('branding.txt', 'r')]
    limit = len(verses)
    track = 0
    for track in range(len(verses)):
        say(verses[track])
        ira.runAndWait()
    end_action()
    '''while True:
        say(verses[track])
        ira.runAndWait()
        track += 1
        if track == limit:
            sleep(15)
            track = 0'''


def start_action():
    try:
        arduino.write("E".encode('ascii'))
        #call_serial("E")
        sleep(0.3)
        # arduino.write(wordlength)
    except AttributeError:
        print("No Arduino is Installed")


def end_action():
    try:
        arduino.write("S".encode('ascii'))
        #call_serial("S")
        sleep(0.3)
    except AttributeError:
        print("No Arduino is Installed")


def say(sentence):
    # word_lengths = [len(word) for word in sentence.split()]
    start_action()
    eye = ["Hi"]
    find_BD = set(sentence.split()).intersection(set(eye))
    if find_BD:
        arduino.write("P".encode('ascii'))
        #call_serial("P")
        sleep(0.3)
    ira.say(sentence)


def match_counter(answer):
    dictionary = [line.rstrip('\n') for line in open('qa.txt', 'r')]
    match_percent = process.extractOne(answer, dictionary)
    ans4match = dictionary.index(match_percent[0])+1
    match_info = list(match_percent)
    match_info[1] = (fuzz.ratio(answer, match_percent[0])+match_percent[1])/2
    match_info.append(dictionary)
    match_info.append(ans4match)
    return match_info


if __name__ == '__main__':
    arduino = arduino_port()
    ira = pyttsx3.init()
    ira.setProperty('rate', 150)
    ira.setProperty('volume', 2)
    #ira.say("hi i am eira")
    ira.say("eira is initializing")
    arduino.write("I".encode('ascii'))
    root = IRACommander()
    bg_file = PhotoImage(file='favicon.gif')
    bg_img = Label(root, image=bg_file)
    bg_img.place(x=0, y=0, relwidth=1, relheight=1)
    tk = CoreGUI(root)
    ira.runAndWait()
    root.mainloop()
